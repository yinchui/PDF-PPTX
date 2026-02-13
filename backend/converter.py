from __future__ import annotations

import io
import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


ProgressCallback = Callable[[int, str, dict[str, Any] | None], None]

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MIN_SHAPE_IN = 0.03


@dataclass
class ConversionOptions:
    mode: str = "local_high_precision"
    vector_tolerance_pt: float = 0.6
    cluster_gap_pt: float = 6.0
    background_filter_ratio: float = 0.35
    min_icon_size_pt: float = 8.0
    max_icon_size_pt: float = 220.0
    debug: bool = False


@dataclass
class JobArtifacts:
    pptx_bytes: bytes
    report: dict[str, Any]
    page_graph: dict[str, Any]


class PdfToPptConverter:
    def __init__(self, options: ConversionOptions):
        self.options = options

    def convert(self, pdf_bytes: bytes, progress: ProgressCallback) -> JobArtifacts:
        progress(5, "开始解析 PDF", None)

        document = fitz.open(stream=pdf_bytes, filetype="pdf")
        total_pages = len(document)

        page_graph: dict[str, Any] = {
            "pages": [],
            "version": "1.0",
            "strategy": "vector-first-with-fallback",
        }
        extracted_pages: list[dict[str, Any]] = []

        report: dict[str, Any] = {
            "total_pages": total_pages,
            "vector_icons_ok": 0,
            "vector_icons_fallback": 0,
            "text_count": 0,
            "image_count": 0,
            "warnings": [],
            "icons": [],
        }

        try:
            for index in range(total_pages):
                page = document[index]
                page_data = self._extract_page(page, index + 1, document)
                extracted_pages.append(page_data)

                report["text_count"] += len(page_data["texts"])
                report["image_count"] += len(page_data["images"])
                page_graph["pages"].append(page_data["page_graph"])

                extract_progress = 10 + int(((index + 1) / max(total_pages, 1)) * 45)
                progress(extract_progress, f"提取对象层（{index + 1}/{total_pages}）", None)

            progress(60, "开始写入 PPTX", None)
            pptx_bytes = self._build_pptx(document, extracted_pages, report, progress)

            report["warnings"] = sorted(set(report["warnings"]))
            page_graph["summary"] = {
                "pages": total_pages,
                "texts": report["text_count"],
                "images": report["image_count"],
                "vector_icons_ok": report["vector_icons_ok"],
                "vector_icons_fallback": report["vector_icons_fallback"],
            }

            progress(100, "转换完成", {"report": report})
            return JobArtifacts(pptx_bytes=pptx_bytes, report=report, page_graph=page_graph)
        finally:
            document.close()

    def _extract_page(self, page: fitz.Page, page_no: int, document: fitz.Document) -> dict[str, Any]:
        page_w = float(page.rect.width)
        page_h = float(page.rect.height)
        page_area = max(1.0, page_w * page_h)

        texts = self._extract_texts(page, page_w, page_h)
        images = self._extract_images(page, document, page_w, page_h)
        vectors = self._extract_vectors(page, page_area)
        icon_candidates = self._build_icon_candidates(vectors, page_w, page_h)

        page_graph = {
            "page_no": page_no,
            "width_pt": page_w,
            "height_pt": page_h,
            "texts": [self._text_for_graph(t) for t in texts],
            "images": [self._image_for_graph(i) for i in images],
            "vectors": [self._vector_for_graph(v) for v in vectors],
            "icons": [self._icon_for_graph(icon) for icon in icon_candidates],
        }

        return {
            "page_no": page_no,
            "page_w": page_w,
            "page_h": page_h,
            "texts": texts,
            "images": images,
            "vectors": vectors,
            "icons": icon_candidates,
            "page_graph": page_graph,
        }

    def _extract_texts(self, page: fitz.Page, page_w: float, page_h: float) -> list[dict[str, Any]]:
        text_dict = page.get_text("dict")
        texts: list[dict[str, Any]] = []

        for block in text_dict.get("blocks", []):
            if block.get("type") != 0:
                continue

            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = str(span.get("text", "")).strip()
                    if not text:
                        continue

                    x0, y0, x1, y1 = _normalize_bbox(span.get("bbox"), page_w, page_h)
                    texts.append(
                        {
                            "text": text,
                            "bbox_pt": (x0, y0, x1, y1),
                            "font_name": str(span.get("font", "Arial")),
                            "font_size_pt": float(span.get("size", 12.0)),
                            "color": _normalize_color(span.get("color")),
                        }
                    )

        texts.sort(key=lambda item: (round(item["bbox_pt"][1], 1), item["bbox_pt"][0]))
        return texts

    def _extract_images(
        self,
        page: fitz.Page,
        document: fitz.Document,
        page_w: float,
        page_h: float,
    ) -> list[dict[str, Any]]:
        images: list[dict[str, Any]] = []
        image_defs = page.get_images(full=True)

        for index, img_def in enumerate(image_defs):
            xref = int(img_def[0])

            try:
                extracted = document.extract_image(xref)
            except Exception:
                continue

            image_bytes = extracted.get("image")
            if not image_bytes:
                continue

            ext = str(extracted.get("ext", "png")).lower()
            mime = f"image/{ext if ext != 'jpg' else 'jpeg'}"
            rects = page.get_image_rects(xref)

            for rect_idx, rect in enumerate(rects):
                x0, y0, x1, y1 = _normalize_bbox(rect, page_w, page_h)
                images.append(
                    {
                        "id": f"img_{xref}_{index}_{rect_idx}",
                        "bbox_pt": (x0, y0, x1, y1),
                        "mime": mime,
                        "bytes": image_bytes,
                    }
                )

        return images

    def _extract_vectors(self, page: fitz.Page, page_area: float) -> list[dict[str, Any]]:
        vectors: list[dict[str, Any]] = []

        for idx, path in enumerate(page.get_drawings()):
            rect = path.get("rect")
            if not rect:
                continue

            bbox = (float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1))
            area = _bbox_area(bbox)
            if area <= 0:
                continue

            if area / page_area > self.options.background_filter_ratio:
                continue

            vectors.append(
                {
                    "id": f"vec_{idx}",
                    "bbox_pt": bbox,
                    "items": path.get("items", []),
                    "stroke": _normalize_color(path.get("color")),
                    "fill": _normalize_color(path.get("fill")),
                    "width": float(path.get("width", 0.75)),
                    "close_path": bool(path.get("closePath", False)),
                    "type": str(path.get("type", "")),
                }
            )

        return vectors

    def _build_icon_candidates(
        self,
        vectors: list[dict[str, Any]],
        page_w: float,
        page_h: float,
    ) -> list[dict[str, Any]]:
        del page_w, page_h
        min_size = self.options.min_icon_size_pt
        max_size = self.options.max_icon_size_pt

        candidates = []
        filtered = []
        for vector in vectors:
            x0, y0, x1, y1 = vector["bbox_pt"]
            width = x1 - x0
            height = y1 - y0
            if width < min_size or height < min_size:
                continue
            if width > max_size or height > max_size:
                continue
            filtered.append(vector)

        if not filtered:
            return candidates

        clusters = _cluster_vectors(filtered, self.options.cluster_gap_pt)
        for cluster_idx, cluster in enumerate(clusters):
            bbox = _union_bbox([item["bbox_pt"] for item in cluster])
            candidates.append(
                {
                    "id": f"icon_{cluster_idx}",
                    "bbox_pt": bbox,
                    "paths": cluster,
                    "classify_result": "vector_candidate",
                }
            )

        return candidates

    def _build_pptx(
        self,
        document: fitz.Document,
        extracted_pages: list[dict[str, Any]],
        report: dict[str, Any],
        progress: ProgressCallback,
    ) -> bytes:
        presentation = Presentation()
        presentation.slide_width = Inches(SLIDE_WIDTH_IN)
        presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
        blank_layout = presentation.slide_layouts[6]

        total_pages = len(extracted_pages)
        for index, page_data in enumerate(extracted_pages):
            slide = presentation.slides.add_slide(blank_layout)
            page_w = page_data["page_w"]
            page_h = page_data["page_h"]

            for text in page_data["texts"]:
                self._add_text(slide, text, page_w, page_h)

            for image in page_data["images"]:
                self._add_image(slide, image["bytes"], image["bbox_pt"], page_w, page_h)

            page = document[page_data["page_no"] - 1]
            for icon in page_data["icons"]:
                icon_record = {
                    "page_no": page_data["page_no"],
                    "icon_id": icon["id"],
                    "bbox_pt": list(icon["bbox_pt"]),
                    "result": "vector",
                    "reason": "",
                }
                try:
                    if not self._add_icon_vector(slide, icon, page_w, page_h):
                        raise ValueError("vector path unsupported")
                    report["vector_icons_ok"] += 1
                except Exception as exc:
                    fallback_bytes = self._rasterize_clip(page, icon["bbox_pt"])
                    self._add_image(slide, fallback_bytes, icon["bbox_pt"], page_w, page_h)
                    report["vector_icons_fallback"] += 1
                    icon_record["result"] = "fallback_image"
                    icon_record["reason"] = str(exc)
                    report["warnings"].append(
                        f"Icon {icon['id']} on page {page_data['page_no']} fallback to image: {exc}"
                    )

                report["icons"].append(icon_record)

            write_progress = 60 + int(((index + 1) / max(total_pages, 1)) * 35)
            progress(write_progress, f"写入幻灯片（{index + 1}/{total_pages}）", None)

        output_stream = io.BytesIO()
        presentation.save(output_stream)
        return output_stream.getvalue()

    def _add_text(self, slide, text: dict[str, Any], page_w: float, page_h: float) -> None:
        x0, y0, x1, y1 = text["bbox_pt"]
        left, top, width, height = _pdf_bbox_to_inches((x0, y0, x1, y1), page_w, page_h)

        box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(max(width, MIN_SHAPE_IN)),
            Inches(max(height, MIN_SHAPE_IN)),
        )
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text["text"]

        scale = (SLIDE_HEIGHT_IN * 72.0) / max(page_h, 1.0)
        font = run.font
        font.size = Pt(max(6.0, min(72.0, text["font_size_pt"] * scale)))
        font.name = text["font_name"] or "Arial"

        rgb = _to_rgb_color(text.get("color"))
        if rgb:
            font.color.rgb = rgb

    def _add_image(
        self,
        slide,
        image_bytes: bytes,
        bbox_pt: tuple[float, float, float, float],
        page_w: float,
        page_h: float,
    ) -> None:
        left, top, width, height = _pdf_bbox_to_inches(bbox_pt, page_w, page_h)
        if width <= 0 or height <= 0:
            return

        slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            Inches(left),
            Inches(top),
            width=Inches(max(width, MIN_SHAPE_IN)),
            height=Inches(max(height, MIN_SHAPE_IN)),
        )

    def _add_icon_vector(self, slide, icon: dict[str, Any], page_w: float, page_h: float) -> bool:
        all_ok = True
        for path in icon["paths"]:
            drawn = self._draw_vector_path(slide, path, page_w, page_h)
            all_ok = all_ok and drawn
        return all_ok

    def _draw_vector_path(self, slide, path: dict[str, Any], page_w: float, page_h: float) -> bool:
        items = path.get("items", [])
        if not items:
            return False

        if len(items) == 1 and _item_op(items[0]) == "re":
            rect = _item_rect(items[0])
            if not rect:
                return False
            self._draw_rectangle(slide, rect, page_w, page_h, path)
            return True

        points, closed = _flatten_path_to_points(items, path.get("close_path", False), self.options.vector_tolerance_pt)
        if len(points) < 2:
            return False

        scale_x = Inches(SLIDE_WIDTH_IN) / max(page_w, 1.0)
        scale_y = Inches(SLIDE_HEIGHT_IN) / max(page_h, 1.0)
        start_x, start_y = points[0]

        builder = slide.shapes.build_freeform(start_x=start_x, start_y=start_y, scale=(scale_x, scale_y))
        builder.add_line_segments(points[1:], close=closed)
        shape = builder.convert_to_shape()
        self._apply_shape_style(shape, path, page_h)
        return True

    def _draw_rectangle(
        self,
        slide,
        rect: tuple[float, float, float, float],
        page_w: float,
        page_h: float,
        path: dict[str, Any],
    ) -> None:
        left, top, width, height = _pdf_bbox_to_inches(rect, page_w, page_h)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(max(width, MIN_SHAPE_IN)),
            Inches(max(height, MIN_SHAPE_IN)),
        )
        self._apply_shape_style(shape, path, page_h)

    def _apply_shape_style(self, shape, path: dict[str, Any], page_h: float) -> None:
        fill_rgb = _to_rgb_color(path.get("fill"))
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        else:
            shape.fill.background()

        stroke_rgb = _to_rgb_color(path.get("stroke"))
        if stroke_rgb:
            shape.line.color.rgb = stroke_rgb
        path_width = max(0.25, float(path.get("width", 0.75)))
        line_scale = (SLIDE_HEIGHT_IN * 72.0) / max(page_h, 1.0)
        shape.line.width = Pt(max(0.25, path_width * line_scale))

    def _rasterize_clip(self, page: fitz.Page, bbox_pt: tuple[float, float, float, float]) -> bytes:
        rect = fitz.Rect(*bbox_pt)
        pix = page.get_pixmap(clip=rect, alpha=True, matrix=fitz.Matrix(2, 2))
        return pix.tobytes("png")

    def _text_for_graph(self, text: dict[str, Any]) -> dict[str, Any]:
        return {
            "text": text["text"],
            "bbox_pt": list(text["bbox_pt"]),
            "font_name": text["font_name"],
            "font_size_pt": text["font_size_pt"],
            "color": text["color"],
        }

    def _image_for_graph(self, image: dict[str, Any]) -> dict[str, Any]:
        return {
            "id": image["id"],
            "bbox_pt": list(image["bbox_pt"]),
            "mime": image["mime"],
        }

    def _vector_for_graph(self, vector: dict[str, Any]) -> dict[str, Any]:
        return {
            "id": vector["id"],
            "bbox_pt": list(vector["bbox_pt"]),
            "stroke": vector["stroke"],
            "fill": vector["fill"],
            "width": vector["width"],
            "type": vector["type"],
            "ops": [_item_op(item) for item in vector["items"]],
        }

    def _icon_for_graph(self, icon: dict[str, Any]) -> dict[str, Any]:
        return {
            "id": icon["id"],
            "bbox_pt": list(icon["bbox_pt"]),
            "paths": [path["id"] for path in icon["paths"]],
            "classify_result": icon["classify_result"],
        }


def _normalize_bbox(
    bbox: Any,
    page_w: float,
    page_h: float,
) -> tuple[float, float, float, float]:
    if isinstance(bbox, fitz.Rect):
        x0, y0, x1, y1 = bbox.x0, bbox.y0, bbox.x1, bbox.y1
    elif isinstance(bbox, (list, tuple)) and len(bbox) >= 4:
        x0, y0, x1, y1 = bbox[0], bbox[1], bbox[2], bbox[3]
    else:
        x0 = y0 = x1 = y1 = 0

    x0 = max(0.0, min(float(x0), page_w))
    y0 = max(0.0, min(float(y0), page_h))
    x1 = max(0.0, min(float(x1), page_w))
    y1 = max(0.0, min(float(y1), page_h))
    if x1 < x0:
        x0, x1 = x1, x0
    if y1 < y0:
        y0, y1 = y1, y0
    return (x0, y0, x1, y1)


def _normalize_color(value: Any) -> list[int] | None:
    rgb = _to_rgb_tuple(value)
    if rgb is None:
        return None
    return [rgb[0], rgb[1], rgb[2]]


def _to_rgb_tuple(value: Any) -> tuple[int, int, int] | None:
    if value is None:
        return None
    if isinstance(value, int):
        return ((value >> 16) & 255, (value >> 8) & 255, value & 255)
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        if all(isinstance(v, (int, float)) for v in value[:3]):
            vals = []
            for v in value[:3]:
                if isinstance(v, float) and 0.0 <= v <= 1.0:
                    vals.append(int(round(v * 255)))
                else:
                    vals.append(int(max(0, min(255, round(v)))))
            return (vals[0], vals[1], vals[2])
    return None


def _to_rgb_color(value: Any) -> RGBColor | None:
    rgb = _to_rgb_tuple(value)
    if rgb is None:
        return None
    return RGBColor(rgb[0], rgb[1], rgb[2])


def _pdf_bbox_to_inches(
    bbox_pt: tuple[float, float, float, float],
    page_w: float,
    page_h: float,
) -> tuple[float, float, float, float]:
    x0, y0, x1, y1 = bbox_pt
    left = (x0 / max(page_w, 1.0)) * SLIDE_WIDTH_IN
    top = (y0 / max(page_h, 1.0)) * SLIDE_HEIGHT_IN
    width = ((x1 - x0) / max(page_w, 1.0)) * SLIDE_WIDTH_IN
    height = ((y1 - y0) / max(page_h, 1.0)) * SLIDE_HEIGHT_IN
    return (left, top, width, height)


def _bbox_area(bbox: tuple[float, float, float, float]) -> float:
    return max(0.0, bbox[2] - bbox[0]) * max(0.0, bbox[3] - bbox[1])


def _union_bbox(bboxes: list[tuple[float, float, float, float]]) -> tuple[float, float, float, float]:
    x0 = min(b[0] for b in bboxes)
    y0 = min(b[1] for b in bboxes)
    x1 = max(b[2] for b in bboxes)
    y1 = max(b[3] for b in bboxes)
    return (x0, y0, x1, y1)


def _cluster_vectors(vectors: list[dict[str, Any]], gap: float) -> list[list[dict[str, Any]]]:
    clusters: list[list[dict[str, Any]]] = []
    visited = [False] * len(vectors)

    for idx in range(len(vectors)):
        if visited[idx]:
            continue
        queue = [idx]
        visited[idx] = True
        cluster = []

        while queue:
            current = queue.pop()
            cluster.append(vectors[current])
            for nxt in range(len(vectors)):
                if visited[nxt]:
                    continue
                if _bbox_close(vectors[current]["bbox_pt"], vectors[nxt]["bbox_pt"], gap):
                    visited[nxt] = True
                    queue.append(nxt)

        clusters.append(cluster)

    return clusters


def _bbox_close(a: tuple[float, float, float, float], b: tuple[float, float, float, float], gap: float) -> bool:
    return not (
        a[2] + gap < b[0]
        or b[2] + gap < a[0]
        or a[3] + gap < b[1]
        or b[3] + gap < a[1]
    )


def _item_op(item: Any) -> str:
    if isinstance(item, (list, tuple)) and item:
        return str(item[0])
    return ""


def _item_rect(item: Any) -> tuple[float, float, float, float] | None:
    if _item_op(item) != "re":
        return None
    if not isinstance(item, (list, tuple)) or len(item) < 2:
        return None
    rect = item[1]
    if isinstance(rect, fitz.Rect):
        return (float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1))
    if isinstance(rect, (list, tuple)) and len(rect) >= 4:
        return (float(rect[0]), float(rect[1]), float(rect[2]), float(rect[3]))
    return None


def _flatten_path_to_points(
    items: list[Any],
    close_path: bool,
    tolerance_pt: float,
) -> tuple[list[tuple[float, float]], bool]:
    points: list[tuple[float, float]] = []

    for item in items:
        op = _item_op(item)
        if op == "re":
            rect = _item_rect(item)
            if not rect:
                continue
            x0, y0, x1, y1 = rect
            rect_points = [(x0, y0), (x1, y0), (x1, y1), (x0, y1)]
            if not points:
                points.extend(rect_points)
            else:
                points.extend(rect_points[1:])
            close_path = True
            continue

        point_args = _extract_point_args(item)
        if op == "l" and len(point_args) >= 2:
            if not points:
                points.append(point_args[0])
            points.append(point_args[1])
            continue

        if op == "c" and len(point_args) >= 4:
            curve = _approximate_cubic_bezier(
                point_args[0],
                point_args[1],
                point_args[2],
                point_args[3],
                tolerance_pt,
            )
            if not points:
                points.append(curve[0])
            points.extend(curve[1:])
            continue

        if op == "m" and point_args:
            if not points:
                points.append(point_args[0])
            continue

        if point_args:
            if not points:
                points.append(point_args[0])
            points.extend(point_args[1:])

    deduped: list[tuple[float, float]] = []
    for p in points:
        if not deduped:
            deduped.append(p)
            continue
        if math.hypot(p[0] - deduped[-1][0], p[1] - deduped[-1][1]) >= 0.01:
            deduped.append(p)

    return deduped, close_path


def _extract_point_args(item: Any) -> list[tuple[float, float]]:
    if not isinstance(item, (list, tuple)):
        return []
    points = []
    for arg in item[1:]:
        if isinstance(arg, fitz.Point):
            points.append((float(arg.x), float(arg.y)))
        elif isinstance(arg, (list, tuple)) and len(arg) >= 2:
            if all(isinstance(v, (int, float)) for v in arg[:2]):
                points.append((float(arg[0]), float(arg[1])))
    return points


def _approximate_cubic_bezier(
    p0: tuple[float, float],
    p1: tuple[float, float],
    p2: tuple[float, float],
    p3: tuple[float, float],
    tolerance: float,
) -> list[tuple[float, float]]:
    chord = math.hypot(p3[0] - p0[0], p3[1] - p0[1])
    segments = max(6, min(30, int(math.ceil(chord / max(tolerance, 0.1)))))
    points = []
    for idx in range(segments + 1):
        t = idx / segments
        mt = 1.0 - t
        x = (
            mt * mt * mt * p0[0]
            + 3 * mt * mt * t * p1[0]
            + 3 * mt * t * t * p2[0]
            + t * t * t * p3[0]
        )
        y = (
            mt * mt * mt * p0[1]
            + 3 * mt * mt * t * p1[1]
            + 3 * mt * t * t * p2[1]
            + t * t * t * p3[1]
        )
        points.append((x, y))
    return points


def write_artifacts(job_dir: Path, artifacts: JobArtifacts) -> tuple[Path, Path, Path]:
    pptx_path = job_dir / "output.pptx"
    report_path = job_dir / "report.json"
    graph_path = job_dir / "page_graph.json"

    pptx_path.write_bytes(artifacts.pptx_bytes)
    report_path.write_text(json.dumps(artifacts.report, ensure_ascii=False, indent=2), encoding="utf-8")
    graph_path.write_text(json.dumps(artifacts.page_graph, ensure_ascii=False, indent=2), encoding="utf-8")

    return pptx_path, report_path, graph_path
